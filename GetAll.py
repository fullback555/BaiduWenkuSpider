# -*- coding: utf-8 -*-
# 爬取百度文库，txt直接下载，ppt下载图片再合成pdf，其他的转成网页再转换为pdf。
# 样式最贴近原文档。
import requests
import chardet
from bs4 import BeautifulSoup
import os
import sys
import shutil  # 文件复制、删除等
import re  # 正则表达式
import json
import time
import math
from PIL import Image
from reportlab.pdfgen import canvas  # pdf生成
import pdfkit  # html to pdf
import pdfcrowd  # html to pdf

from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from config import MyPath

# 多线程与多进程
from concurrent.futures import ThreadPoolExecutor, as_completed
from concurrent.futures import ProcessPoolExecutor

from PyQt5.QtCore import QObject, pyqtSignal, pyqtSlot, QRunnable


class WorkerSignals(QObject):
	"""定义信号"""
	finished = pyqtSignal(str)  # 任务完成信号
	condition = pyqtSignal(str)  # 当前状态信号
	# result = pyqtSignal(object)  # 结果
	progress = pyqtSignal(int)  # 进度


class GetAll(QRunnable):

	def __init__(self, url, savepath=None):
		"""
		:param url: 待爬取文档所在页面的url
		:param savepath: 生成文档保存路径
		"""
		super(GetAll, self).__init__()
		self.signals = WorkerSignals()

		self.headers = {
			"User-Agent": "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/65.0.3325.146 Safari/537.36"
		}
		self.savepath0 = savepath
		# 百度文库doc类型超过50页会放在不同的url中
		self.kv = None  # 传参到url
		self.url = url
		self.startpage = 1  # 开始页数，构造网页地址使用
		self.imagepath = ''  # 存放图片路径
		self.mode = 'w'  # txt的写入模式默认为覆盖写入
		# self.init()

	def init(self):
		self.html = self.get_html_text(self.url, self.kv)
		self.html2 = None  # BeautifullSoup格式化过的网页
		# 存储文档基本信息:title,docType,docId,totalPageNum
		wkinfo = self.getWkInfo()
		self.wkinfo = wkinfo if wkinfo else self.getVipWkInfo()

		# 设置默认下载路径
		self.savepath = self.makeDirs(self.savepath0, self.wkinfo.get('docType'))  # 默认保存路径
		self.htmlpath = self.makeDirs(self.savepath0, 'html')
		self.htmlfile = self.wkinfo.get('title')+".html"

		# txt/ppt无此类url
		if self.wkinfo.get('docType') in ['doc', 'xls', 'pdf'] and self.wkinfo.get('is_vip_free_doc'): 
			self.html2 = self.getSoupHtml(self.url, self.kv)  # BeautifullSoup格式化过的网页
			self.jsonurls = self.getJsonUrl()  # doc文档json的url和图片url
		# 否则获取付费文档的
		elif self.wkinfo.get('docType') in ['doc', 'xls', 'pdf'] and not self.wkinfo.get('is_vip_free_doc'):
			self.html2 = self.getSoupHtml(self.url, self.kv)  # BeautifullSoup格式化过的网页
			self.jsonurls = self.getVipJsonUrls()  # doc文档json/png的url
		

	# 获取当前程序路径
	def get_cur_path(self, filename=None):
		path_dir = os.path.dirname(os.path.realpath(sys.argv[0]))
		if not filename:  # 如果为空
			return path_dir
		path = os.path.join(path_dir, filename)
		return path

	# 创建文件夹,在path路径下再创建文档类型命名的文件夹
	def makeDirs(self, path, docType):
		if not docType:  # 如果docType类型为空
			docType = 'NoType'
		if os.path.exists(path):  # 如果path是绝对路径并已存在
			temppath = os.path.join(path, docType)
		else:
			temppath = os.path.join(self.get_cur_path(path), docType)

		if not os.path.exists(temppath):
			os.makedirs(temppath)
			return temppath
		return temppath

	# 获取网站源代码，使用BeautifulSoup解析
	def getSoupHtml(self, url, kv=None):
		session = requests.session()
		try:
			r = session.get(url, headers=self.headers, params=kv)
			# 根据检测到网页的编码方式来解码
			r.encoding = chardet.detect(r.content).get('encoding')
			html = BeautifulSoup(r.text, 'html.parser')  # 格式化html
			return html
		except Exception as e:
			print(f'获取网页格式化出错：{e}')
			return None

	# 获取网站源代码，返回文本内容。kv是键值对，可以在url后面添加参数，如&pn=51
	def get_html_text(self, url, kv=None):
		session = requests.session()
		try:
			r = session.get(url, headers=self.headers, params=kv)
			# 根据检测到网页的编码方式来解码
			r.encoding = chardet.detect(r.content).get('encoding')
			# r.encoding = r.apparent_encoding  # 或者
			# self.html = r.text
			print(r.url)  # 打印url
			return r.text
			# return BeautifulSoup(r.text, 'html.parser')  # 格式化
		except Exception as e:
			print(f'获取网页出错：{e}')
			return None

	# 获取网站源代码，返回二进制内容
	def get_html_content(self, url):
		session = requests.session()
		try:
			r = session.get(url, headers=self.headers)
			# 根据检测到网页的编码方式来解码
			# r.encoding = chardet.detect(r.content).get('encoding')
			return r.content
			# return BeautifulSoup(r.text, 'html.parser')  # 格式化
		except Exception as e:
			print(f'获取网页出错：{e}')
			return None

	# 获取文档基本信息：名称、文档ID、类型、页数。常见文档可用此法
	def getWkInfo(self):
		infoDict = {}  # 用于保存文档基本信息
		items = ['title', 'docId', 'docType', 'totalPageNum', 'is_vip_free_doc']
		for key in items:
			item = re.findall(f"\'{key}\':.*?\'(.*?)\',", str(self.html))  # 无？号，最大匹配原则
			if item:
				infoDict[key] = str(item[0])
				if key == 'is_vip_free_doc':  # 是否是vip免费文档，以此判断是否是付费文档
					if item[0] == 'false':
						infoDict[key] = False
					else:
						infoDict[key] = True
		print('文档基本信息：', infoDict)
		return infoDict

	# 获取文档基本信息：名称、文档ID、类型、页数。付费文档
	def getVipWkInfo(self):
		infoDict = {}  # 用于保存文档基本信息
		items = ['title', 'docId', 'docType', 'totalPageNum', 'is_vip_free_doc']
		vip_items = ['title', 'show_doc_id', 'type', 'page', 'is_vip_free_doc']
		# 文档格式号与文档格式对应表
		type_dict = {'8': 'txt', '4': 'doc', '5': 'xls', '3': 'ppt', '7': 'pdf', '1': 'doc', '6': 'ppt'}
		for i, key in enumerate(vip_items):
			item = re.findall(f'\"{key}\":.*?\"?(.*?)\"?,', str(self.html))  # 最大匹配原则
			if item:
				infoDict[items[i]] = str(item[0])
				if items[i] == 'docType':  # 将文档格式号转换为文档格式
					infoDict[items[i]] = type_dict.get(str(item[0]))
				if items[i] == 'is_vip_free_doc':  # 是否是vip免费文档，以此判断是否是付费文档
					if item[0] == 'false':
						infoDict[items[i]] = False
					else:
						infoDict[items[i]] = True
		print('VIP文档基本信息：', infoDict)
		return infoDict

	# 从url中获取json，并返回字典
	def getJson(self, url):
		'''url: json文件所在页面的url'''
		html = self.get_html_text(url)
		# 获取json格式数据
		if str(html)[-1] == ')':
			jsonStr = re.findall(r"\((.*)\)", str(html))[0]
		else:
			jsonStr = str(html)
		# print('json:', jsonStr)
		# 将json字符串解析为python的对象
		textdict = json.loads(jsonStr)
		# print('json to dict:', textdict)
		return textdict

	# 解析获取txt文档，或者doc文档的文本部分。
	def parse_txt(self):
		json_url = 'https://wenku.baidu.com/api/doc/getdocinfo?callback=cb&doc_id=' +\
			self.wkinfo.get('docId')
		print(json_url)
		textdict = self.getJson(json_url)
		# print(textdict)

		# json文件所在url的参数
		content_url = 'https://wkretype.bdimg.com/retype/text/' + \
			textdict.get('doc_id') + \
			textdict.get('md5sum') + \
			'&pn=1&rn=' + \
			textdict.get('docInfo').get('totalPageNum') + \
			'&type=txt' + '&rsign=' + textdict.get('rsign')
		print(content_url)
		# 解析json，转换成python对象
		content = self.getJson(content_url)
		print(content)
		# 发射进度信号
		self.signals.progress.emit(10)

		string = ''
		for j, item in enumerate(content):  # 每页一个parags
			for i in item['parags']:
				string += i['c'].replace('\\r', '\r').replace('\\n', '\n')
			# 发射进度信号
			self.signals.progress.emit(60*((j+1)/len(content)) + 10)
		return string

	# 保存txt文档内容，content:要写入的内容, mode:写入模式，如：'w', 'a'
	def saveToTxt(self, content, mode='w'):
		savepath = os.path.join(self.savepath, self.wkinfo.get('title') + '.txt')
		with open(savepath, mode, encoding='utf-8') as f:
			f.write(content)
		return savepath


	# 获取txt文档
	def get_txt(self):
		content = self.parse_txt()
		txt_file = self.saveToTxt(content)
		return txt_file

	# 解析获取ppt/pdf文档，返回包含图片的url链接
	def parse_ppt(self, docId=None):
		if docId is None:
			docId = self.wkinfo.get('docId')
		# 文档信息的链接
		json_url = 'https://wenku.baidu.com/browse/getbcsurl?doc_id=' + \
			docId + '&pn=1&rn=99999&type=ppt'
		print(json_url)
		# 解析json
		img_json = self.getJson(json_url)
		# print(img_json)
		# 获取ppt图片的url地址
		picurls = [item.get('zoom') for item in img_json]
		picurls = self.filter_picurls(picurls)  # 过滤假的picurls
		print('ppt picurls:', picurls)
		return picurls

	# 过滤假的picurls。如'&png=63974-74826&jpg=1256626-1373523'
	def filter_picurls(self, picurls):
		if not picurls:
			return []
		urls = list(filter(self.check_picurl, picurls))
		return urls

	def check_picurl(self, picurl):
		jpg = picurl.split('&jpg=')[-1]
		start = jpg.split('-')[0]
		end = jpg.split('-')[-1]
		if start == end or not start or not end:
			return False
		return True

	# 下载单个图片，url:图片链接，filename:要保存到本地的路径
	def download_pic(self, url, filename):
		if not url:
			return None  # 如果url为空/None就跳过
		content = self.get_html_content(url)
		# 保存为图片
		with open(filename, 'wb') as f:
			f.write(content)
			# print(f'已下载图片：{filename}')
		return filename


	# 下载图片, 参数urls为图片url地址，默认保存为png格式，返回下载后图片路径。使用线程池。
	def download_pics(self, urls):  # 下载图片，并存储
		if not urls:  
			return []
		# 创建临时保存图片文件夹
		imagepath = os.path.join('images', self.wkinfo.get('title'))
		self.imagepath = self.makeDirs(self.savepath, imagepath)
		picpaths = []  # 储存图片的名字和路径，在合成ppt时保持正确的顺序

		# 创建10个线程的线程池
		with ThreadPoolExecutor(max_workers=10) as t:
			threads = []
			for index, url in enumerate(urls):
				if not url: continue  # 如果url为空/None就跳过
				filename = os.path.join(self.imagepath, str(index + 1) + '.png')
				picpaths.append(filename)
				# 下载图片
				# self.download_pic(url, filename)
				obj = t.submit(self.download_pic, url, filename)
				threads.append(obj)

			i = 1
			# as_completed方法是一个生成器，没任务完成时会一直阻塞，除非设置了timeout
			# 当有任务完成时，会yield这个任务，先完成的任务会先返回给主线程。先执行完的先返回结果
			for future in as_completed(threads):
				data = future.result()
				print(f'成功下载图片: {data}')
				# 发射进度信号
				self.signals.progress.emit(int(70*i/len(urls)) + 10)
				i += 1 

		print(f'图片保存在：{self.imagepath}')
		return picpaths  # 返回正确顺序的图片路径和名字

	# 根据下载的图片生成ppt
	def pic_to_ppt(self, picpaths):
		if not picpaths:
			return None
		# 如果是ppt， 在上面下载好图片后，将图片重新合成ppt
		prs = Presentation()  # 空白ppt文档
		for picpath in picpaths:
			# 一张幻灯片
			slide = prs.slides.add_slide(prs.slide_layouts[6])
			left, top, width, height = Inches(0), Inches(0), Inches(10), Inches(7.5)
			pic = slide.shapes.add_picture(picpath, left, top, width, height)
		# 保存ppt
		prs.save(os.path.join(self.savepath, self.wkinfo.get('title') + '.pptx'))
		# 发射进度信号
		self.signals.progress.emit(85)
		return os.path.join(self.savepath, self.wkinfo.get('title') + '.pptx')

	# 根据下载的图片生成pdf，生成pdf的速度比较快。
	def pic_to_pdf(self, picpaths):	
		if not picpaths:
			return None
		img = Image.open(picpaths[0])
		w, h = img.size[:2]

		pdf = os.path.join(self.savepath, self.wkinfo.get('title') + '.pdf')
		c = canvas.Canvas(pdf, pagesize=(w, h))  # 默认为A4大小
		for picpath in picpaths:
			try:
				img = Image.open(picpath)
			except Exception as e:
				print('pic_to_pdf，打开图片错误原因：{}'.format(e))
				continue
			# img = Image.open(picpath)
			w, h = img.size[:2]
			c.drawImage(picpath, 0, 0, w, h)
			c.showPage()  # 保存当前画布页面
		c.save()  # 保存文件并关闭画布
		# 发射进度信号
		self.signals.progress.emit(100)
		return pdf
		
	# 清除下载的图片
	def remove_file(self, path):
		# 删除文件夹
		if os.path.exists(path) and os.path.isdir(path):
			shutil.rmtree(path)
		# 删除文件
		elif os.path.exists(path) and os.path.isfile(path):
			os.remove(path)

	# 获取ppt文档，并清除下载的图片
	def get_ppt(self):
		picurls = self.parse_ppt()
		# 发射进度信号
		self.signals.progress.emit(10)
		picpaths = self.download_pics(picurls)
		ppt_file = self.pic_to_ppt(picpaths)
		print(f'已生成文件：{ppt_file}')
		pdf_file = self.pic_to_pdf(picpaths)
		print(f'已生成文件：{pdf_file}')
		self.signals.progress.emit(100)
		# 清除下载的图片
		# self.remove_file(self.imagepath)
		# print('临时图片文件夹imgae已删除')
		return ppt_file


	# 将获取的图片合成pdf文件
	# def mergeImageToPDF(self, pages):
	# 	if pages == 0:
	# 		raise IOError

	# 	namelist = [os.path.join(self.tempdirpath, str(x)+'.png') for x in range(pages)]
	# 	firstimg = Image.open(namelist[0])
	# 	imglist = []
	# 	for i, imgname in enumerate(namelist[1:]):
	# 		img = Image.open(imgname)
	# 		img.load()

	# 		if img.mode == 'RGBA':  # png图片的转为RGB mode,否则保存时会引发异常
	# 			img.mode = 'RGB'
	# 		imglist.append(img)
	# 		# 发射进度信号
	# 		self.signals.progress.emit(int(40*(i+1)/len(namelist[1:]))+50)

	# 	savepath = os.path.join(self.pptsavepath, self.wkinfo.get('title')+'.pdf')
	# 	firstimg.save(savepath, "PDF", resolution=100.0,
	# 				  save_all=True, append_images=imglist)
	# 	return savepath


	# 获取存储信息的json文件的url，包括0.json/0.png
	def getJsonUrl(self):
		urlinfo = re.findall(r"WkInfo.htmlUrls = \'(.*?)\';", str(self.html))
		if urlinfo:
			pass
			urlinfo = str(urlinfo[0]).replace('\\', '').replace('x22', '')
			allurls = re.findall(r"(https:.*?)}", urlinfo)
			print(allurls)
			self.jsonurls = allurls
			return allurls
		return []

	# 获取付费文档的json、png文件的url
	def getVipJsonUrls(self): 
		jsonurls = re.findall('pageLoadUrl.*?(https:.*?0.json.*?)\\\\\"}', self.html)
		jsonurls = [addr.replace("\\\\/", "/") for addr in jsonurls]

		picurls = []
		pngs = re.findall('\"png.*?(pageLoadUrl.*?https:.*0.png.*\\\\\"})', self.html)
		if pngs:
			picurls = re.findall('pageLoadUrl.*?(https:.*?0.png.*?)\\\\\"}', pngs[0])
			picurls = [addr.replace("\\\\/", "/") for addr in picurls]
		
		self.jsonurls = jsonurls + picurls	
		print('jsonurls:', self.jsonurls)
		return self.jsonurls

	# 判断文档是否为ppt格式
	def isPptStyle(self):
		# iswholepic = False
		ispptlike = False
		for url in self.jsonurls:
			if "0.json" in url:
				textdict = self.getJson(url)
				# 若json文件中的style属性为空字符串且font属性为None,则说明pdf全由图片组成
				# if textdict.get("style") == "" and textdict.get("font") is None:
				# 	iswholepic = True
				# 	break
				if textdict.get('page').get('pptlike'):
					ispptlike = True
			break
		# return iswholepic and ispptlike
		return ispptlike

	# 创建html文档,用于组织爬取的文件
	def creatHtml(self):
		with open(os.path.join(self.htmlpath, str(self.startpage) + self.htmlfile), "w") as f:
			# 生成文档头
			message = """
			<!DOCTYPE html>
			<html class="expanded screen-max">
				<head>
				<meta charset="utf-8">
				<title>文库</title>"""
			f.write(message)


	def addMessageToHtml(self, message):
		""":param message:向html文档中添加内容 """
		with open(os.path.join(self.htmlpath, str(self.startpage) + self.htmlfile), "a", encoding='utf-8') as a:
			a.write(message)


	# 从html中匹配出与控制格式相关的css文件的url
	def getCssUrl(self):
		pattern =  re.compile('<link href="//.*?\.css')
		allmessage =  pattern.findall(str(self.html2))
		allcss = [x.replace('<link href="', "https:") for x in allmessage]
		return allcss

	def getPageTag(self):
		print('html2:', self.html2)
		""":return:返回id属性包含 data-page-no 的所有标签,即所有页面的首标签"""
		def attributeFilter(tag):
			return tag.has_attr('data-page-no')
		return self.html2.find_all(attributeFilter)

	def getDocIdUpdate(self):
		""":return:doc_id_update字符串"""
		pattern = re.compile('doc_id_update:".*?"')
		for i in pattern.findall(str(self.html2)):
			return i.split('"')[1]

	def getAllReaderRenderStyle(self):
		""":return: style <id = "reader - render - style">全部内容"""
		page = 1
		style = '<style id='+'"reader-render-style">\n'
		for url in self.jsonurls:
			if "json" in url:
				textdict = self.getJson(url)
				style += self.getReaderRenderStyle(textdict.get('style'), textdict.get('font'), textdict.get('font'), page)
				page += 1
			else:
				break
		style += "</style>\n"

		return style

	def getReaderRenderStyle(self, allstyle, font, r, page):
		"""
		:param allstyle: json数据的style内容
		:param font: json数据的font内容
		:param r: TODO:解析作用未知,先取值与e相同
		:param page: 当前页面
		:return: style <id = "reader - render - style">
		"""
		p, stylecontent = "", []
		for index in range(len(allstyle)):
			style = allstyle[index]
			if style.get('s'):
				p = self.getPartReaderRenderStyle(style.get('s'), font, r).strip(" ")
			l = "reader-word-s" + str(page) + "-"
			p and stylecontent.append("." + l + (",." + l).join([str(x) for x in style.get('c')]) + "{ " + p + "}")
			if style.get('s').get("font-family"):
				pass
		stylecontent.append("#pageNo-" + str(page) + " .reader-parent{visibility:visible;}")
		return "".join(stylecontent)

	def getPartReaderRenderStyle(self, s, font, r):
		"""
		:param s:  json style下的s属性
		:param font:  json font属性
		:param r: fontMapping TODO:先取为与e相同
		:return: style <id = "reader - render - style">中的部分字符串
		"""
		content = []
		n, p = 10, 1349.19 / 1262.85  # n为倍数, p为比例系数, 通过页面宽度比得出

		def fontsize(f):
			content.append("font-size:" + str(math.floor(eval(f) * n * p)) + "px;")

		def letterspacing(l):
			content.append("letter-spacing:" + str(eval(l) * n) + "px;")

		def bold(b):
			"false" == b or content.append("font-weight:600;")

		def fontfamily(o):
			n = font.get(o) or o if font else o
			content.append("font-family:'" + n + "','" + o + "','" + (r.get(n) and r[n] or n) + "';")

		for attribute in s:
			if attribute == "font-size":
				fontsize(s[attribute])
			elif attribute == "letter-spacing":
				letterspacing(s[attribute])
			elif attribute == "bold":
				bold(s[attribute])
			elif attribute == "font-family":
				fontfamily(s[attribute])
			else:
				content.append(attribute + ":" + s[attribute] + ";")
		return "".join(content)

	# 向html中添加css
	def AddCss(self):
		urls = self.getCssUrl()

		print('css urls:', urls)
		urls = [url  for url in urls if "htmlReader" in url or "core" in url or "main" in url or "base" in url]
		for i, url in enumerate(urls):
			message = '<style type="text/css">'+requests.get(url).text+"</style>>"
			self.addMessageToHtml(message)
			# 发射进度信号
			self.signals.progress.emit(int(20*(i+1)/len(urls)))
		print('addCss urls:', len(urls))

		content = self.getAllReaderRenderStyle()  # 获取文本控制属性css
		self.addMessageToHtml(content)

	def addMainContent(self):
		"""
		:param startpage: 开始生成的页面数
		:return:
		"""
		self.addMessageToHtml("\n\n\n<body>\n")
		docidupdate = self.getDocIdUpdate()

		# 分别获取json和png所在的url
		jsonurl = [x for x in self.jsonurls if "json" in x]
		pngurl = [x for x in self.jsonurls if "png" in x]

		tags = self.getPageTag()

		print('*'*50)
		print(tags)

		for page, tag in enumerate(tags):
			if page > 50:
				break
			tag['style'] = "height: 1349.19px;"
			tag['id'] = "pageNo-" + str(page+1)
			self.addMessageToHtml(str(tag).replace('</div>', ''))
			diu = self.getDocIdUpdate()
			n = "-webkit-transform:scale(1.00);-webkit-transform-origin:left top;"
			textdict = self.getJson(jsonurl[page])

			# 判断是否出现图片url少于json文件url情况
			if page < len(pngurl):
				maincontent = self.creatMainContent(textdict.get('body'), textdict.get('page'), textdict.get('font'), page + 1, docidupdate,
													pngurl[page])
			else:
				maincontent = self.creatMainContent(textdict.get('body'), textdict.get('page'), textdict.get('font'), page + 1, docidupdate, "")
			content = "".join([
				'<div class="reader-parent-' + diu + " reader-parent " + '" style="position:relative;top:0;left:0;' + n + '">',
				'<div class="reader-wrap' + diu + '" style="position:absolute;top:0;left:0;width:100%;height:100%;">',
				'<div class="reader-main-' + diu + '" style="position:relative;top:0;left:0;width:100%;height:100%;">', maincontent,
				"</div>", "</div>", "</div>", "</div>"])

			self.addMessageToHtml(content)
			print("已完成%s页的写入,当前写入进度为%f" % (str(page+self.startpage), 100*(page+self.startpage)/int(self.wkinfo.get('totalPageNum'))) + '%')

			# 发射进度信号
			self.signals.progress.emit(int( 40*(page+self.startpage)/int(self.wkinfo.get('totalPageNum')) + 20))
		
		self.addMessageToHtml("\n\n\n</body>\n</html>")

	# 判断是否为数字
	def isNumber(self, obj):
		return isinstance(obj, int) or isinstance(obj, float)

	def creatMainContent(self, body, page, font, currentpage, o, pngurl):
		"""
		:param body: body属性
		:param page: page属性
		:param font: font属性
		:param currentpage: 当前页面数
		:param o:doc_id_update
		:param pngurl: 图片所在url
		:return:文本及图片的html内容字符串
		"""
		content, p, s, h = 0, 0, 0, 0
		main = []
		l = 2
		c = page.get('v')

		d = font   # d原本为fongmapping
		y = {
				"pic": '<div class="reader-pic-layer" style="z-index:__NUM__"><div class="ie-fix">',
				"word": '<div class="reader-txt-layer" style="z-index:__NUM__"><div class="ie-fix">'
			}
		g = "</div></div>"
		MAX1 , MAX2 = 0, 0
		body = sorted(body, key=lambda k: k.get('p').get('z'))
		for index in range(len(body)):
			content = body[index]
			if "pic" == content.get('t'):
				MAX1 = max(MAX1, content.get('c').get('ih') + content.get('c').get('iy') + 5)
				MAX2 = max(MAX2, content.get('c').get('iw'))
		for index in range(len(body)):
			content = body[index]
			s = content.get('t')
			if not p:
				p = h = s
			if p == s:
				if content.get('t') == "word":
					# m函数需要接受可变参数
					main.append(self.creatTagOfWord(content, currentpage, font, d, c))
				elif content.get('t') == 'pic':
					main.append(self.creatTagOfImage(content, pngurl, MAX1, MAX2))
			else:
				main.append(g)
				main.append(y.get(s).replace('__NUM__', str(l)))
				l += 1
				if content.get('t') == "word":
					# m函数需要接受可变参数
					main.append(self.creatTagOfWord(content, currentpage, font, d, c))
				elif content.get('t') == 'pic':
					main.append(self.creatTagOfImage(content, pngurl, MAX1, MAX2))
				p = s
		return y.get(h).replace('__NUM__', "1") + "".join(main) + g

	def creatTagOfWord(self, t, currentpage, font, o, version, *args):
		"""
		:param t: body中的每个属性
		:param currentpage: page
		:param font: font属性
		:param o:font属性
		:param version: page中的version属性
		:param args:
		:return:<p>标签--文本内容
		"""
		p = t.get('p')
		ps = t.get('ps')
		s = t.get('s')
		z = ['<b style="font-family:simsun;">&nbsp</b>', "\n"]
		k, N = 10, 1349.19 / 1262.85
		# T = self.j
		U = self.O(ps)
		w, h, y, x, D= p.get('w'), p.get('h'), p.get('y'), p.get('x'), p.get('z')
		pattern=re.compile("[\s\t\0xa0]| [\0xa0\s\t]$")
		final = []

		if U and ps and ((ps.get('_opacity') and ps.get('_opacity') == 1) or (ps.get('_alpha') and ps.get('_alpha') == 0)):
			return ""
		else:
			width = math.floor(w * k * N)
			height = math.floor(h * k * N)
			final.append("<p "+'class="'+"reader-word-layer" + self.processStyleOfR(t.get('r'), currentpage) + '" ' + 'style="' + "width:" +str(width) + "px;" + "height:" + str(height) + "px;" + "line-height:" + str(height) + "px;")
			final.append("top:"+str(math.floor(y * k * N))+"px;"+"left:"+str(math.floor(x * k * N))+"px;"+"z-index:"+str(D)+";")
			final.append(self.processStyleOfS(s, font, o, version))
			final.append(self.processStyleOf_rotate(ps.get('_rotate'), w, h, x, y, k, N) if U and ps and self.isNumber(ps.get('_rotate')) else "")
			final.append(self.processStyleOfOpacity(ps.get('_opacity')) if U and ps and ps.get('_opacity') else "")
			final.append(self.processStyleOf_scaleX(ps.get('_scaleX'), width, height) if U and ps and ps.get('_scaleX') else "")
			final.append(str(isinstance(t.get('c'), str) and len(t.get('c')) == 1 and pattern.match(t.get('c')) and "font-family:simsun;") if isinstance(t.get('c'), str) and len(t.get('c')) == 1 and pattern.match(t.get('c')) else "")
			final.append('">')
			final.append(t.get('c') if t.get('c') else "")
			final.append(U and ps and str(self.isNumber(ps.get('_enter'))) and z[ps.get('_enter') if ps.get('_enter') else 1] or "")
			final.append("</p>")

			return "".join(final)

	def processStyleOfS(self, t, font, r, version):
		"""
		:param t: 文本的s属性
		:param font: font属性
		:param r:font属性
		:param version:
		:return:处理好的S属性字符串
		"""
		infoOfS = []
		n = {"font-size": 1}
		p , u = 10, 1349.19 / 1262.85

		def fontfamily(o):
			n = font.get(o) or o if font else o
			if abs(version) > 5:
				infoOfS.append("font-family:'"+ n + "','" + o + "','" + (r.get('n') and r[n] or n) + "';")
			else:
				infoOfS.append("font-family:'" + o + "','" + n + "','" + (r.get(n) and r[n] or n) + "';")

		def bold(e):
			"false" == e or infoOfS.append("font-weight:600;")

		def letter(e):
			infoOfS.append("letter-spacing:" + str(eval(e) * p) + "px;")

		if t is not None:
			for attribute in t:
				if attribute == "font-family":
					fontfamily(t[attribute])
				elif attribute == "bold":
					bold(t[attribute])
				elif attribute == "letter-spacing":
					letter(t[attribute])
				else:
					infoOfS.append(attribute + ":" + (str(math.floor(((t[attribute] if self.isNumber(t[attribute]) else eval(t[attribute])) * p * u))) + "px" if n.get(attribute) else t[attribute]) + ";")

		return "".join(infoOfS)

	def processStyleOfR(self, r, page):
		"""
		:param r: 文本的r属性
		:param page: 当前页面
		:return:
		"""
		l = " " + "reader-word-s" + str(page) + "-"
		return "".join([l + str(x) for x in r]) if isinstance(r, list) and len(r) != 0 else ""

	def processStyleOf_rotate(self, t, w, h, x, y, k, N):
		"""
		:param t: _rotate属性
		:param w: body中p.w
		:param h: body中p.h
		:param x: body中p.x
		:param y: body中p.y
		:param k: 倍数10
		:param N: 比例系数
		:return: 处理好的_rotate属性字符串
		"""
		p = []
		s = k * N
		if t == 90:
			p.append("left:" + str(math.floor(x + (w - h) / 2) * s) + "px;" + "top:" + str(math.floor(y - (h - w) / 2) * s) + "px;" + "text-align: right;" + "height:" + str(math.floor(h + 7) * s) + "px;")
		elif t == 180:
			p.append("left:" + str(math.floor(x - w) * s) + "px;" + "top:" + str(math.floor(y - h) * s) + "px;")
		elif t == 270:
			p.append("left:" + str(math.floor(x + (h - w) / 2) * s) + "px;" + "top:" + str(math.floor(y - (w - h) / 2) * s) + "px;")

		return "-webkit-"+"transform:rotate("+str(t)+"deg);"+"".join(p)

	def processStyleOf_scaleX(self, t, width, height):
		"""
		:param t:     _scaleX属性
		:param width: 计算好的页面width
		:param height:计算好的页面height
		:return: 处理好的_scaleX属性字符串
		"""
		return "-webkit-" + "transform: scaleX(" + str(t) + ");" + "-webkit-" + "transform-origin:left top;width:" + str(width + math.floor(width / 2)) + "px;height:" + str(height + math.floor(height / 2)) + "px;"

	def processStyleOfOpacity(self,t):
		"""
		:param t: opacity属性
		:return:处理好的opacity属性字符串
		"""
		t = (t or 0),
		return "opacity:" + str(t) + ";"

	def creatTagOfImage(self,t,url, *args):
		"""
		:param t: 图片的字典
		:param url:图片链接
		:param args:
		:return:图像标签
		"""
		u, l = t.get('p'), t.get('c')
		if u.get("opacity") and u.get('opacity') == 0:
			return ""
		else:
			if u.get("x1") or (u.get('rotate') != 0 and u.get('opacity') != 1):
				message = '<div class="reader-pic-item" style="' + "background-image: url(" + url + ");" + "background-position:" + str(-l.get('ix')) + "px " + str(-l.get('iy')) + "px;" \
						  + "width:" + str(l.get('iw')) + "px;" + "height:" + str(l.get('ih')) + "px;" + self.getStyleOfImage(u, l) + 'position:absolute;overflow:hidden;"></div>'
			else:
				[s, h] = [str(x) for x in args]
				message = '<p class="reader-pic-item" style="' + "width:" + str(l.get('iw')) + "px;" + "height:" + str(l.get('ih')) + "px;" + self.getStyleOfImage(u, l) + 'position:absolute;overflow:hidden;"><img width="' + str(h) + '" height="' + str(s) + '" style="position:absolute;top:-' + str(l.get('iy')) + "px;left:-" + str(l.get('ix')) + "px;clip:rect(" + str(l.get('iy')) + "px," + str(int(h) - l.get('ix')) + "px, " + str(s) + "px, " + str(l.get('ix')) + 'px);" src="' + url + '" alt=""></p>'

			return message

	def getStyleOfImage(self, t, e):
		"""
		:param t: 图片p属性
		:param e: 图片c属性
		:return:
		"""
		def parseFloat(string):
			"""
			:param string:待处理的字符串
			:return: 返回字符串中的首个有效float值，若字符首位为非数字，则返回nan
			"""
			if string is None:
				return math.nan
			elif isinstance(string, float):
				return string
			elif isinstance(string, int):
				return float(string)
			elif string[0] != ' ' and not str.isdigit(string[0]):
				return math.nan
			else:
				p = re.compile("\d+\.?\d*")
				all = p.findall(string)
				return float(all[0]) if len(all) != 0 else math.nan

		if t is None:
			return ""
		else:
			r, o, a, n = 0, 0, "", 0
			iw = e.get('iw')
			ih = e.get('ih')
			u = 1349.19 / 1262.85
			l = str(t.get('x') * u) + "px"
			c = str(t.get('y') * u) + "px"
			d = ""
			x = {}
			w = {"opacity": 1, "rotate": 1, "z": 1}
			for n in t:
				x[n] = t[n] * u if (self.isNumber(t[n]) and not w.get(n)) else t[n]

			if x.get('w') != iw or x.get('h') != ih:
				if x.get('x1'):
					a = self.P(x.get('x0'), x.get('y0'), x.get('x1'), x.get('y1'), x.get('x2'), x.get('y2'))
				r = parseFloat(parseFloat(a[0])/iw if len(a) else x.get('w') / iw)
				o = parseFloat(parseFloat(a[1])/ih if len(a) else x.get('h') / ih)

				m, v = iw * (r-1), ih * (o-1)
				c = str((x.get('y1') + x.get('y3')) / 2 - parseFloat(ih) / 2)+"px" if x.get('x1') else str(x.get('y') + v / 2) + "px"
				l = str((x.get('x1') + x.get('x3')) / 2 - parseFloat(iw) / 2)+"px" if x.get('x1') else str(x.get('x') + m / 2) + "px"
				d = "-webkit-" + "transform:scale(" + str(r) + "," + str(o) + ")"

			message = "z-index:" + str(x.get('z')) + ";" + "left:" + l + ";" + "top:" + c + ";" + "opacity:" + str(x.get('opacity') or 1) + ";"
			if x.get('x1'):
				message += self.O(x.get('rotate')) if x.get('rotate') > 0.01 else self.O(0, x.get('x1'), x.get('x2'), x.get('y0'), x.get('y1'), d)
			else:
				message += d + ";"

			return message

	def P(self,t, e, r, i, o, a):
		p = round(math.sqrt(math.pow(abs(t - r), 2) + math.pow(abs(e - i), 2)), 4)
		s = round(math.sqrt(math.pow(abs(r - o), 2) + math.pow(abs(i - a), 2)), 4)
		return [s, p]

	def O(self, t, *args):
		[e, r, i, o, a] = [0, 0, 0, 0, ""] if len(args) == 0 else [x for x in args]
		n = o > i
		p = e > r
		if n and p:
			a += " Matrix(1,0,0,-1,0,0)"
		elif n:
			a += " Matrix(1,0,0,-1,0,0)"
		elif p:
			a += " Matrix(-1,0,0,1,0,0)"
		elif t:
			a += " rotate(" + str(t) + "deg)"
		return a + ";"

	# 将网页转换为pdf
	def convertHtmlToPdf(self):
		savepath = os.path.join(self.savepath, str(self.startpage)+self.wkinfo.get('title') + '.pdf')

		# 每个url的最大页数为50
		exactpages = int(self.wkinfo.get('totalPageNum'))
		if exactpages > 50:
			exactpages = 50

		# 如果资源文件目录中存在wkhtmltopdf.app，则复制到应用程序目录中来。操作系统不同，此处需修改。
		path_wk = self.get_cur_path('wkhtmltopdf')
		if not os.path.exists(path_wk):
			QMessageBox.information(self, '提示', '请将文件wkhtmltopdf拷贝到wenku.app/Contents/MacOS')
			return False
			# MyPath.copy_file(filename='wkhtmltopdf')
		# (需下载wkhtmltox)将程序路径传入config对象
		config = pdfkit.configuration(wkhtmltopdf=path_wk)

		options = {'disable-smart-shrinking': '',
				   'lowquality': '',
				   'image-quality': 60,
				   'page-height': str(1349.19*0.26458333),
				   'page-width': '291',
				   'margin-bottom': '0',
				   'margin-top': '0',
				   }
		pdfkit.from_file(os.path.join(self.htmlpath, str(self.startpage) + self.htmlfile), savepath, options=options, configuration=config)

		return savepath

	# 使用pdfcrowd把html转化为pdf，需要用户名和key，试用每天30次。
	def htmlToPdf(self):
		savepath = os.path.join(self.savepath, str(self.startpage)+self.wkinfo.get('title') + '.pdf')

		try:
			# create the API client instance
			client = pdfcrowd.HtmlToPdfClient('your_name', 'your_key')

			# run the conversion and write the result to a file
			# client.convertUrlToFile('http://www.example.com', 'example.pdf')
			client.setNoMargins(True)  # 边距为0, 默认页面大小为A4
			client.convertFileToFile(os.path.join(self.htmlpath, str(self.startpage) + self.htmlfile), savepath)
		except pdfcrowd.Error as why:
			# report the error
			print('Pdfcrowd Error: {}\n'.format(why))
			return False
			# sys.stderr.write('Pdfcrowd Error: {}\n'.format(why))

		return savepath


	# 开始线程，该函数自动运行
	@pyqtSlot()
	def run(self):
		self.init()
		docType = self.wkinfo.get('docType')
		filename = '{}.{}'.format(self.wkinfo.get('title'), self.wkinfo.get('docType'))

		# 判断文档是否为txt格式
		if docType == 'txt':
			self.signals.condition.emit('当前下载文档：{}'.format(filename))
			savepath = self.get_txt()
			# 发射进度信号
			self.signals.progress.emit(100)
			self.signals.finished.emit(str(savepath))
		
		# 判断文档是否为ppt格式
		elif docType == 'ppt':
			self.signals.condition.emit('当前下载文档：{}'.format(filename))
			savepath = self.get_ppt()
			self.signals.finished.emit(str(savepath))

		# doc/xls/pdf下载
		elif docType in ['doc', 'xls', 'pdf']:
			# 类似于ppt的pdf，使用ppt的获取方法，超过50页也不影响。
			if docType == 'pdf' and self.isPptStyle():
				self.signals.condition.emit('当前下载文档：{}'.format(filename))
				savepath = self.get_ppt()
				self.signals.finished.emit(str(savepath))
				return
			if not self.html2:  # 如果获取网页失败则退出
				self.signals.finished.emit(str(False))  # 发射失败信号
				return

			for page in range(math.ceil(int(self.wkinfo.get('totalPageNum'))/50)):
				self.startpage = page * 50 + 1
				self.signals.condition.emit('当前下载文档：{}'.format(str(self.startpage) + filename))
				if page == 0:
					self.creatHtml()  # 创建html文档,用于组织爬取的文件
					self.AddCss()  # 向html中添加css
					self.addMainContent()  # 生成页面
					savepath = self.convertHtmlToPdf()  # 将html转换成pdf
					if not savepath:  # pdfcrowd转换方法失败后
						savepath = self.htmlToPdf()  # 将html转换成pdf

					# 发射进度信号
					self.signals.progress.emit(100)

				else:
					self.kv = {'pn': self.startpage}
					self.init()

					self.creatHtml()  # 创建html文档,用于组织爬取的文件
					self.AddCss()  # 向html中添加css
					self.addMainContent()  # 生成页面
					savepath = self.convertHtmlToPdf()  # 将html转换成pdf
					if not savepath:  # pdfcrowd转换方法失败后
						savepath = self.htmlToPdf()  # 将html转换成pdf

					# 发射进度信号
					self.signals.progress.emit(100)
			self.signals.finished.emit(str(savepath))  # 发射完成信号

		else:
			print('{}该类型文档无法解析'.format(docType))


if __name__ == '__main__':
	pass
	# 若存储路径为空，则在当前文件夹生成
	# GetAll('https://wenku.baidu.com/view/fb92d7d3b8d528ea81c758f5f61fb7360a4c2b61.html?from=search',
	# 			  "文库").run()

