# -*- coding: utf-8 -*-
# @Author: wenlong
# @Date:   2020-05-04 16:31:02
# @Last Modified by:   fullback555
# @Last Modified time: 2020-05-29 23:26:26

# 爬取百度文库文档，付费文档只能下载试读部分内容。

import requests
import chardet  # 检测网页编码
# from bs4 import BeautifulSoup  # 解析
import os
import sys
import shutil  # 文件复制、删除等
import re  # 正则表达式
import json
import time
import math
from PIL import Image
from reportlab.pdfgen import canvas  # pdf生成
# from docx import Document
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
# 个人写的类
from my_docx import Mydocx
from my_image import MyImage
from config import MyPath
# 多线程与多进程
from concurrent.futures import ThreadPoolExecutor, as_completed
from concurrent.futures import ProcessPoolExecutor
# from threading import Thread  # 多线程
# from multiprocessing import Process, Pool  # 多进程

from PyQt5.QtCore import QObject, pyqtSignal, pyqtSlot, QRunnable


class WorkerSignals(QObject):
	"""定义信号"""
	finished = pyqtSignal(str)  # 任务完成信号
	condition = pyqtSignal(str)  # 当前状态信号
	# result = pyqtSignal(object)  # 结果
	progress = pyqtSignal(int)  # 进度


class WenKuSpider(QRunnable):
	"""百度文库爬虫，获取百度文库文档。"""
	def __init__(self, url, savepath):
		"""
		:param url: 待爬取文档所在页面的url
		:param savepath: 生成文档保存路径
		"""
		super(WenKuSpider, self).__init__()
		self.signals = WorkerSignals()

		self.headers = {
			"User-Agent": "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/65.0.3325.146 Safari/537.36"
		}
		self.savepath0 = savepath
		# 百度文库doc类型超过50页会放在不同的url中
		self.kv = None  # 传参到url
		self.url = url
		self.startpage = 1  # 开始页数，构造网页地址使用
		self.imagepath = ''  # 存放图片路径
		self.mode = 'w'  # txt的写入模式默认为覆盖写入
		self.isPppImage = False  # 是否ppt图片
		
		# self.init()

	# 初始化
	def init(self):
		self.html = self.get_html_text(self.url, self.kv)
		# 存储文档基本信息:title,docType,docId,totalPageNum
		wkinfo = self.getWkInfo()
		self.wkinfo = wkinfo if wkinfo else self.getVipWkInfo()
		self.totalpage = int(self.wkinfo.get('totalPageNum', '0'))  # 获取总页数
		# print(self.wkinfo)
		# 设置默认下载路径
		self.savepath = self.makeDirs(self.savepath0, self.wkinfo.get('docType'))  # 默认保存路径
		# txt/ppt无此类url
		if self.wkinfo.get('docType') in ['doc', 'xls', 'pdf'] and self.wkinfo.get('is_vip_free_doc'): 
			self.allurls = self.getAllUrls()  # doc文档json的url和图片url
			self.jsonurls = self.getJsonUrls(self.allurls)  # doc文档json的url, 注意要在getPicUrls前执行
			self.picurls = self.getPicUrls(self.allurls)  # doc文档图片的url
		# 否则获取付费文档的
		elif self.wkinfo.get('docType') in ['doc', 'xls', 'pdf'] and not self.wkinfo.get('is_vip_free_doc'):
			self.jsonurls = self.getVipJsonUrls()  # doc文档json的url
			self.picurls = self.getVipPicUrls()  # doc文档图片的url

	# 获取当前程序路径
	def get_cur_path(self, filename=None):
		path_dir = os.path.dirname(os.path.realpath(sys.argv[0]))
		if not filename:  # 如果为空
			return path_dir
		path = os.path.join(path_dir, filename)
		return path


	# 创建文件夹,在path路径下再创建文档类型命名的文件夹
	def makeDirs(self, path, docType):
		if not docType:  # 如果docType类型为空
			docType = 'NoType'
		if os.path.exists(path):  # 如果path是绝对路径并已存在
			temppath = os.path.join(path, docType)
		else:
			temppath = os.path.join(self.get_cur_path(path), docType)

		if not os.path.exists(temppath):
			os.makedirs(temppath)
			return temppath
		return temppath

	# 获取网站源代码，返回文本内容。kv是键值对，可以在url后面添加参数，如&pn=51
	def get_html_text(self, url, kv=None):
		session = requests.session()
		try:
			r = session.get(url, headers=self.headers, params=kv)
			# 根据检测到网页的编码方式来解码
			r.encoding = chardet.detect(r.content).get('encoding')
			# r.encoding = r.apparent_encoding  # 或者
			# self.html = r.text
			print(r.url)  # 打印url
			return r.text
			# return BeautifulSoup(r.text, 'html.parser')  # 格式化
		except Exception as e:
			print(f'获取网页出错：{e}')
			return None

	# 获取网站源代码，返回二进制内容
	def get_html_content(self, url):
		session = requests.session()
		try:
			r = session.get(url, headers=self.headers)
			# 根据检测到网页的编码方式来解码
			# r.encoding = chardet.detect(r.content).get('encoding')
			return r.content
			# return BeautifulSoup(r.text, 'html.parser')  # 格式化
		except Exception as e:
			print(f'获取网页出错：{e}')
			return None

	# 获取文档基本信息：名称、文档ID、类型、页数。常见文档可用此法
	def getWkInfo(self):
		infoDict = {}  # 用于保存文档基本信息
		items = ['title', 'docId', 'docType', 'totalPageNum', 'is_vip_free_doc']
		for key in items:
			item = re.findall(f"\'{key}\':.*?\'(.*?)\',", str(self.html))  # 无？号，最大匹配原则
			if item:
				infoDict[key] = str(item[0])
				if key == 'is_vip_free_doc':  # 是否是vip免费文档，以此判断是否是付费文档
					if item[0] == 'false':
						infoDict[key] = False
					else:
						infoDict[key] = True
		print('文档基本信息：', infoDict)
		return infoDict

	# 获取文档基本信息：名称、文档ID、类型、页数。付费文档
	def getVipWkInfo(self):
		infoDict = {}  # 用于保存文档基本信息
		items = ['title', 'docId', 'docType', 'totalPageNum', 'is_vip_free_doc']
		vip_items = ['title', 'show_doc_id', 'type', 'page', 'is_vip_free_doc']
		# 文档格式号与文档格式对应表
		type_dict = {'8': 'txt', '4': 'doc', '5': 'xls', '3': 'ppt', '7': 'pdf', '1': 'doc', '6': 'ppt'}
		for i, key in enumerate(vip_items):
			item = re.findall(f'\"{key}\":.*?\"?(.*?)\"?,', str(self.html))  # 最大匹配原则
			if item:
				infoDict[items[i]] = str(item[0])
				if items[i] == 'docType':  # 将文档格式号转换为文档格式
					infoDict[items[i]] = type_dict.get(str(item[0]))
				if items[i] == 'is_vip_free_doc':  # 是否是vip免费文档，以此判断是否是付费文档
					if item[0] == 'false':
						infoDict[items[i]] = False
					else:
						infoDict[items[i]] = True
		print('VIP文档基本信息：', infoDict)
		return infoDict

	# 从url中获取json，并返回字典
	def getJson(self, url):
		'''url: json文件所在页面的url'''
		html = self.get_html_text(url)
		# 获取json格式数据
		if str(html)[-1] == ')':
			jsonStr = re.findall(r"\((.*)\)", str(html))[0]
		else:
			jsonStr = str(html)
		# print('json:', jsonStr)
		# 将json字符串解析为python的对象
		try:
			textdict = json.loads(jsonStr)
		except Exception as e:
			print('解析json异常')
			textdict = None
		# print('json to dict:', textdict)
		return textdict

	# 解析获取txt文档，或者doc文档的文本部分。
	def parse_txt(self):
		json_url = 'https://wenku.baidu.com/api/doc/getdocinfo?callback=cb&doc_id=' +\
			self.wkinfo.get('docId')
		print(json_url)
		textdict = self.getJson(json_url)
		# print(textdict)

		# json文件所在url的参数
		content_url = 'https://wkretype.bdimg.com/retype/text/' + \
			textdict.get('doc_id') + \
			textdict.get('md5sum') + \
			'&pn=1&rn=' + \
			textdict.get('docInfo').get('totalPageNum') + \
			'&type=txt' + '&rsign=' + textdict.get('rsign')
		print(content_url)
		# 解析json，转换成python对象
		content = self.getJson(content_url)
		if not content:  # 如果没有获取到内容
			return ''
		# 发射进度信号
		# self.signals.progress.emit(5)

		string = ''
		for j, item in enumerate(content):  # 每页一个parags
			for i in item['parags']:
				string += i['c'].replace('\\r', '\r').replace('\\n', '\n')
			# 发射进度信号
			self.signals.progress.emit(50*((j+1)/len(content)) + 10)
		return string

	# 保存txt文档内容，content:要写入的内容, mode:写入模式，如：'w', 'a'
	def saveToTxt(self, content, mode='w'):
		savepath = os.path.join(self.savepath, self.wkinfo.get('title') + '.txt')
		with open(savepath, mode, encoding='utf-8') as f:
			f.write(content)
		return savepath

	# 获取txt文档
	def get_txt(self):
		content = self.parse_txt()
		txt_file = self.saveToTxt(content)
		return txt_file


	# 解析获取ppt/pdf文档，返回包含图片的url链接
	def parse_ppt(self, docId=None):
		if docId is None:
			docId = self.wkinfo.get('docId')
		# 文档信息的链接
		json_url = 'https://wenku.baidu.com/browse/getbcsurl?doc_id=' + \
			docId + '&pn=1&rn=99999&type=ppt'
		print(json_url)
		# 解析json
		img_json = self.getJson(json_url)
		# print(img_json)
		# 获取ppt图片的url地址
		picurls = [item.get('zoom') for item in img_json]
		picurls = self.filter_picurls(picurls)  # 过滤假的picurls
		print('ppt picurls:', picurls)
		return picurls

	# 过滤假的picurls。如'&png=63974-74826&jpg=1256626-1373523'
	def filter_picurls(self, picurls):
		if not picurls:
			return []
		urls = list(filter(self.check_picurl, picurls))
		return urls

	def check_picurl(self, picurl):
		jpg = picurl.split('&jpg=')[-1]
		start = jpg.split('-')[0]
		end = jpg.split('-')[-1]
		if start == end or not start or not end:
			return False
		return True


	# 下载单个图片，url:图片链接，filename:要保存到本地的路径
	def download_pic(self, url, filename):
		if not url:
			return None  # 如果url为空/None就跳过
		content = self.get_html_content(url)
		# 保存为图片
		with open(filename, 'wb') as f:
			f.write(content)
			# print(f'已下载图片：{filename}')
		return filename

	# 下载图片, 参数urls为图片url地址，默认保存为png格式，返回下载后图片路径。没有使用多线程。
	# def download_pics(self, urls):  # 下载图片，并存储
	# 	# 创建临时保存图片文件夹
	# 	imagepath = os.path.join('images', self.wkinfo.get('title'))
	# 	self.imagepath = self.makeDirs(self.savepath, imagepath)
	# 	picpaths = []  # 储存图片的名字和路径，在合成ppt时保持正确的顺序
	# 	for index, url in enumerate(urls):
	# 		if not url: continue  # 如果url为空/None就跳过
	# 		filename = os.path.join(self.imagepath, str(index + 1) + '.png')
	# 		picpaths.append(filename)
	# 		# 下载图片
	# 		self.download_pic(url, filename)
	# 	print(f'图片保存在：{self.imagepath}')
	# 	return picpaths  # 返回正确顺序的图片路径和名字

	# 下载图片, 参数urls为图片url地址，默认保存为png格式，返回下载后图片路径。使用线程池。
	def download_pics(self, urls):  # 下载图片，并存储
		if not urls:  
			return []
		# 创建临时保存图片文件夹
		imagepath = os.path.join('images', self.wkinfo.get('title'))
		self.imagepath = self.makeDirs(self.savepath, imagepath)
		picpaths = []  # 储存图片的名字和路径，在合成ppt时保持正确的顺序

		# 创建10个线程的线程池
		with ThreadPoolExecutor(max_workers=10) as t:
			threads = []
			for index, url in enumerate(urls):
				if not url: continue  # 如果url为空/None就跳过
				index += self.startpage
				filename = os.path.join(self.imagepath, str(index) + '.png')
				picpaths.append(filename)
				# 下载图片
				# self.download_pic(url, filename)
				obj = t.submit(self.download_pic, url, filename)
				threads.append(obj)

			i = 1
			# as_completed方法是一个生成器，没任务完成时会一直阻塞，除非设置了timeout
			# 当有任务完成时，会yield这个任务，先完成的任务会先返回给主线程。先执行完的先返回结果
			for future in as_completed(threads):
				data = future.result()
				print(f'成功下载图片: {data}')
				# 发射进度信号
				if self.isPppImage:  # 如果下载的ppt图片
					self.signals.progress.emit(int(70*i/len(urls)) + 10)
				else:
					self.signals.progress.emit(int(25*i/len(urls)) + 25)
				i += 1 

		print(f'图片保存在：{self.imagepath}')
		return picpaths  # 返回正确顺序的图片路径和名字

	# 根据下载的图片生成ppt
	def pic_to_ppt(self, picpaths):
		if not picpaths:
			return None
		# 如果资源文件目录中存在模板default.pptx，则复制到应用程序目录中来。
		filename = self.get_cur_path('default.pptx')
		if not os.path.exists(filename):
			MyPath.copy_file(filename='default.pptx')
		# 如果是ppt， 在上面下载好图片后，将图片重新合成ppt
		try:
			prs = Presentation(filename)  # 空白ppt文档
		except Exception as e:
			prs = Presentation()  # 空白ppt文档
		
		for picpath in picpaths:
			# 一张幻灯片
			slide = prs.slides.add_slide(prs.slide_layouts[6])
			left, top, width, height = Inches(0), Inches(0), Inches(10), Inches(7.5)
			pic = slide.shapes.add_picture(picpath, left, top, width, height)
		# 保存ppt
		prs.save(os.path.join(self.savepath, self.wkinfo.get('title') + '.pptx'))
		return os.path.join(self.savepath, self.wkinfo.get('title') + '.pptx')

	# 根据下载的图片生成pdf，生成pdf的速度比较快。
	def pic_to_pdf(self, picpaths):	
		if not picpaths:
			return None
		img = Image.open(picpaths[0])
		w, h = img.size[:2]

		pdf = os.path.join(self.savepath, self.wkinfo.get('title') + '.pdf')
		c = canvas.Canvas(pdf, pagesize=(w, h))  # 默认为A4大小
		for picpath in picpaths:
			try:
				img = Image.open(picpath)
			except Exception as e:
				print('pic_to_pdf，打开图片错误原因：{}'.format(e))
				continue
			# img = Image.open(picpath)
			w, h = img.size[:2]
			c.drawImage(picpath, 0, 0, w, h)
			c.showPage()  # 保存当前画布页面
		c.save()  # 保存文件并关闭画布
		return pdf
		
	# 清除下载的图片
	def remove_file(self, path):
		# 删除文件夹
		if os.path.exists(path) and os.path.isdir(path):
			shutil.rmtree(path)
		# 删除文件
		elif os.path.exists(path) and os.path.isfile(path):
			os.remove(path)

	# 获取ppt文档，并清除下载的图片
	def get_ppt(self):
		# 获取文字内容
		txt = self.parse_txt()
		self.saveToTxt(txt)
		# 获取图片，转换成ppt或pdf
		picurls = self.parse_ppt()
		# 发射进度信号
		self.signals.progress.emit(10)
		self.isPppImage = True  # 是ppt图片
		picpaths = self.download_pics(picurls)
		ppt_file = self.pic_to_ppt(picpaths)
		print(f'已生成文件：{ppt_file}')
		pdf_file = self.pic_to_pdf(picpaths)
		print(f'已生成文件：{pdf_file}')
		self.signals.progress.emit(100)
		# 清除下载的图片
		# self.remove_file(self.imagepath)
		# print('临时图片文件夹imgae已删除')
		return ppt_file

	# 获取pdf文档
	def get_pdf(self):
		print('pdf isPppLike: ', self.isPptStyle())
		# 如果类似于ppt，则使用ppt方法解析下载
		if self.isPptStyle():
			# 获取文字内容
			txt = self.parse_txt()
			self.saveToTxt(txt)
			# 获取图片，转换成pdf
			picurls = self.parse_ppt()
			# 发射进度信号
			self.signals.progress.emit(10)
			self.isPppImage = True  # 是ppt图片
			picpaths = self.download_pics(picurls)
			pdf_file = self.pic_to_pdf(picpaths)
			print(f'已生成文件：{pdf_file}')
			self.signals.progress.emit(100)
			return pdf_file
		else:  # 如果与ppt不类似，需要另行解决
			savepath = self.get_doc()
			return savepath
			# pdf_file = self.pic_to_pdf(picpaths)
			# print(f'已生成文件：{pdf_file}')
			# return pdf_file

	# 获取所有的json文件的url，包括0.json和0.png
	def getAllUrls(self):
		urlinfo = re.findall(r"WkInfo.htmlUrls = \'(.*?)\';", str(self.html))[0]
		urlinfo = str(urlinfo).replace('\\', '').replace('x22', '')
		# print(urlinfo)
		allurls = re.findall(r"(https:.*?)}", urlinfo)
		# allurls = [str(url) for url in allurls]  # 所有的urls
		return allurls

	# 获取存储信息的json文件的url
	def getJsonUrls(self, allurls):
		jsonurls = [str(url) for url in allurls if '0.json' in url]  # 包含有json数据的urls
		print('jsonurls no:', len(jsonurls))
		return jsonurls

	# 获取存储信息的json图片pic的url
	def getPicUrls(self, allurls):
		picurls = [str(url) for url in allurls if '0.png' in url]  # 包含图片的urls
		# 获取有图片的页码
		pages_pic = self.get_pages_has_pic()
		# 如果该页无图片，则将url设置为None，不作删除是为了让哪页上图片对应在哪页上。
		for i, picurl in enumerate(picurls):
			if i+1 not in pages_pic:
				picurls[i] = None
		print('picurls no:', len(picurls))
		return picurls

	# 获取付费文档存储文字信息的json文件的url
	def getVipJsonUrls(self): 
		jsonurls = re.findall('pageLoadUrl.*?(https:.*?0.json.*?)\\\\\"}', self.html)
		jsonurls = [addr.replace("\\\\/", "/") for addr in jsonurls]
		print('vip jsonurls:', jsonurls)
		return jsonurls

	# 获取付费文档存储图片信息的png文件的url
	def getVipPicUrls(self):
		pngs = re.findall('\"png.*?(pageLoadUrl.*?https:.*0.png.*\\\\\"})', self.html)
		if pngs:
			# print(pngs[0])
			picurls = re.findall('pageLoadUrl.*?(https:.*?0.png.*?)\\\\\"}', pngs[0])
			picurls = [addr.replace("\\\\/", "/") for addr in picurls]
			print('vip picurls:', picurls)
			return picurls
		return []

	# 获取doc文档哪一页有图片，返回有图片的页码列表
	def get_pages_has_pic(self):
		# self.jsonurls = self.getJsonUrls()
		pages_pic = []
		for i, url in enumerate(self.jsonurls):  # url链接数等于页数
			pic_flag = False
			textdict = self.getJson(url)
			# 如果textdict没有'body'这个键值，则跳过。下面几个if判断是避免出错
			if 'body' not in textdict.keys():
				continue
			if type(textdict['body']) != list:
				continue
			for item in textdict['body']:
				# 图片
				if not item: continue
				if 't' not in item.keys() or 's' not in item.keys():
					continue
				if item['t'] == "pic":
					if item['s'] is not None:
						pic_flag = True
			if pic_flag:
				pages_pic.append(i+1)
		return pages_pic
		
	# 切割图片，去除图片空白部分，返回切割后图片路径，没有使用多进程
	# def cut_image(self, picpaths):
	# 	paths = []
	# 	# 切割图片，去除多余的空白部分
	# 	for picpath in picpaths:
	# 		pic = MyImage(picpath)
	# 		paths += pic.cut_image()
	# 	print('picpaths no:', len(paths), paths)
	# 	return paths 

	# 切割图片，去除图片空白部分，返回切割后图片路径。应用多进程提高速度（方式1）。
	# def cut_image(self, picpaths):
	# 	paths = []
	# 	# 创建4个进程，进程数量与cpu核数量相当即可。
	# 	pool = Pool(processes=4)
	# 	results = []
	# 	# 切割图片，去除多余的空白部分
	# 	for picpath in picpaths:
	# 		pic = MyImage(picpath)
	# 		# apply_async()是apply()的并行版本，本身就可以返回被进程调用的函数的返回值
	# 		results.append(pool.apply_async(pic.cut_image, ()))  # 返回值的对象（不是值本身）		
	# 	pool.close()  # 关闭进程池，表示不能再往进程池中添加进程，需要在join之前调用
	# 	pool.join()  # 等待进程池中的所有进程执行完毕
	# 	# 获取返回值, 注意结果是有序的。
	# 	for res in results:
	# 		paths += res.get()
	# 	print('picpaths no:', len(paths), paths)
	# 	return paths 

	# 切割图片，去除图片空白部分，返回切割后图片路径。应用多进程提高速度（方式2）。
	def cut_image(self, picpaths):
		paths = []
		# 不指定max_workers，则默认是cpu的个数
		p = ProcessPoolExecutor()
		threads = []
		# 切割图片，去除多余的空白部分
		for picpath in picpaths:
			pic = MyImage(picpath)
			obj = p.submit(pic.cut_image)
			threads.append(obj)
		# 发射进度信号
		self.signals.progress.emit(60)
		# 类似multiprocessing的close和join一起使用，不写这句，就用with上下文管理
		p.shutdown() 
		# 获取返回结果，结果是有序的，按原提交函数顺利返回结果。
		for future in threads:
			paths += future.result()

		# 发射进度信号
		self.signals.progress.emit(70)

		print('picpaths no:', len(paths), paths)
		return paths 

	# 解析获取doc文档的文字部分，方法不同于parse_txt，效果有时略好。
	def parse_doc(self): 
		result = ''  # 用于保存文本
		jsonurls = re.findall('(https.*?0.json.*?)\\\\x22}', self.html)
		jsonurls = [addr.replace("\\\\\\/", "/") for addr in jsonurls]
		print(jsonurls)
		for url in jsonurls:
			textdict = self.getJson(url)
			y = 0  # 文档中的y坐标，用于判断是否有换行。
			# 如果textdict没有'body'这个键值，则跳过
			if 'body' not in textdict.keys():
				continue
			if type(textdict['body']) != list:
				continue
			for item in textdict['body']:
				# word文本
				if item['t'] == "word":
					if not y == item['p'].get('y'):   #一行都有一个标记值，可用标记值来判断是否元素在同一行
						y = item['p'].get('y')   
						n = '\n'           #若是y与标记值不相等，等则行，用'\n'
					else:
						n = ''             #若是y与标记值相等，则不换行，用''
					result += n            #连接时先连接n,再连接文字，这样可以很好的处理换行。
					# result += item['c'].encode('utf-8').decode('unicode_escape', 'ignore')
					result += item['c']
		return result

	# 解析doc文档的文字和图片，参数paths为下载图片，切割去空白后图片路径
	def parse_doc2(self, paths):
		filename = os.path.join(self.savepath, self.wkinfo.get('title') + '.docx')
		doc = Mydocx(filename)
		index = 0
		for i, url in enumerate(self.jsonurls):  # url链接数等于页数
			# 处理文字
			textdict = self.getJson(url)
			# 如果textdict没有'body'这个键值，则跳过
			if 'body' not in textdict.keys():
				continue
			if type(textdict['body']) != list:
				continue

			y = 0  # 文档中的y坐标，用于判断是否有换行。
			word = ''  # 用于保存文本
			pic_flag = False  # 本页是否有图片标志
			for item in textdict['body']:
				# word文本
				if item['t'] == "word":
					if not y == item['p'].get('y'):   #一行都有一个标记值，可用标记值来判断是否元素在同一行
						y = item['p'].get('y')   
						n = '\n'           #若是y与标记值不相等，等则行，用'\n'
					else:
						n = ''             #若是y与标记值相等，则不换行，用''
					word += n            #连接时先连接n,再连接文字，这样可以很好的处理换行。
					# word += item['c'].encode('utf-8').decode('unicode_escape', 'ignore')
					word += item['c']

				# 图片
				elif item['t'] == "pic":
					# print('aaaaaaaaa我是图片aaaaaaaaaaaa')
					if item['s'] is not None:
						pic_flag = True  # 有图片
						pic_url = item['s'].get('pic_file').replace('\\/', '/')
						# pic_no.append(re.findall(r'_(.*)\.', pic_url)[0])
						pic_w = item['p'].get('w')/1.5  # 获取图片宽度，此处修正。
						print('{}, w:{}'.format(pic_url, pic_w))

						# 此处插入文字，图片之前的文字
						doc.add_para()
						doc.add_text(word)
						word = ''
						# print('index:{}'.format(index))
						if index >= len(paths):
							continue
						picname = os.path.basename(paths[index])
						# 如果当前要插入图片与图片列表不符，则插入下一张图。局部修正
						if picname.split('_')[0] != re.findall(r'_(.*)_', pic_url)[0] and index < len(paths)-1:
							index += 1
						# 此处插入图片
						doc.add_enter()  # 换行
						doc.add_pic_in_run(filename=paths[index], width=pic_w/28.346)  # 72dpi 1cm = 28.346像素
						print(f'index:{index}, pic:{paths[index]}')
						index += 1

			print('word:', word)
			# 发射进度信号
			self.signals.progress.emit(int(25*(i+1)/len(self.jsonurls)) + 75)

			# 插入文字，图片之后的文字
			doc.add_para()
			doc.add_text(word)
		
		doc.save(filename)  # 如果原路径文件已存在，则内容追回在后面。
		return filename

	# 获取doc文档的文字和图片，有图片的文档建议使用此方法
	def get_doc(self):
		# 首先下载doc文档中的图片
		self.isPppImage = False  # 不是ppt图片
		picpaths = self.download_pics(self.picurls)
		# 分割图片，去除空白
		paths = self.cut_image(picpaths)
		# 获取doc文档
		savepath = self.parse_doc2(paths)
		return savepath

	# 获取xls文档
	def get_xls(self, mode='w'):
		# 只获取里面的文字
		result = self.parse_txt()  # 速度较快，但一些转义字符不能正确处理，如\/
		# result = self.parse_doc()    # 速度可能略慢，效果较好。
		txtfile = self.saveToTxt(result, mode)
		if mode == 'w':  # 以下只需要执行一遍，超过50页也没关系。
			# 获取图片合成pdf格式
			picurls = self.parse_ppt()
			self.isPppImage = True  # 类似ppt图片
			picpaths = self.download_pics(picurls)
			pdf = self.pic_to_pdf(picpaths)
			return pdf
		return txtfile

	# 判断文档是否为类似ppt格式(主要以图片展示)
	def isPptStyle(self):
		ispptlike = False
		for url in self.jsonurls:
			textdict = self.getJson(url)
			# print(textdict)
			if textdict.get('page').get('pptlike'):
				ispptlike = True
			break
		return ispptlike

	
	# 开始线程，该函数自动运行，运行获取文档。
	@pyqtSlot()
	def run(self):
		# 初始化
		self.init()
		docType = self.wkinfo.get('docType')
		filename = '{}.{}'.format(self.wkinfo.get('title'), self.wkinfo.get('docType'))
		
		# doc超过50页，网址有变化，txt/ppt/xls/pdf则没有变化
		if docType == 'txt':
			self.signals.progress.emit(10)
			self.signals.condition.emit('当前下载文档：{}'.format(filename))
			savepath = self.get_txt()
			# 发射进度信号
			self.signals.progress.emit(100)
			self.signals.finished.emit(savepath)

		elif docType == 'ppt':
			self.signals.condition.emit('当前下载文档：{}'.format(filename))
			savepath = self.get_ppt()
			self.signals.finished.emit(savepath)

		elif docType == 'pdf':
			self.signals.condition.emit('当前下载文档：{}'.format(filename))
			savepath = self.get_pdf()
			self.signals.finished.emit(savepath)

		elif docType == 'xls':
			for page in range(math.ceil(self.totalpage/50)):
				self.startpage = page * 50 + 1

				# 发射当前状态
				endpage = self.totalpage if self.totalpage<=(page+1)*50 else (page+1)*50
				filename = '{}({}).{}'.format(self.wkinfo.get('title'), str(self.startpage)+'-'+str(endpage)+'页', 'pdf')
				self.signals.condition.emit('当前下载文档：{}'.format(filename))

				self.signals.progress.emit(10)  # 发射进度信号

				if page != 0:  # 超过50页，网址有变化。
					# self.url = re.sub(r'&pn=(.*)', '&pn=' + str(self.startpage), self.url)
					self.mode = 'a'  # 追加写入txt
					self.kv = {'pn': self.startpage}
					self.init()
				savepath = self.get_xls(mode=self.mode)
				self.signals.progress.emit(100)  # 发射进度信号
			self.signals.finished.emit(savepath)

		elif docType == 'doc':
			# 如果存在先移除
			filename = os.path.join(self.savepath, self.wkinfo.get('title') + '.docx')
			if os.path.exists(filename):
				os.remove(filename)
			for page in range(math.ceil(self.totalpage/50)):
				
				self.startpage = page * 50 + 1
				# 发射当前状态
				endpage = self.totalpage if self.totalpage<=(page+1)*50 else (page+1)*50
				filename = '{}({}).{}'.format(self.wkinfo.get('title'), str(self.startpage)+'-'+str(endpage)+'页', 'docx')
				self.signals.condition.emit('当前下载文档：{}'.format(filename))

				self.signals.progress.emit(25)  # 发射进度信号

				if page != 0:  # 超过50页，网址有变化。
					# self.url = re.sub(r'&pn=(.*)', '&pn=' + str(self.startpage), self.url)
					self.kv = {'pn': self.startpage}
					self.init()
				
				savepath = self.get_doc()
			self.signals.finished.emit(savepath)
		else:
			print('{}该类型文档无法解析'.format(docType))
		

# 主函数
if __name__ == '__main__':

	txt_url = 'https://wenku.baidu.com/view/df3abfc36137ee06eff9183f.html?from=search'
	# ppt_url = 'https://wenku.baidu.com/view/a5fc216dc9d376eeaeaad1f34693daef5ff7130b.html?from=search'
	pdf_url = 'https://wenku.baidu.com/view/cdd05642c67da26925c52cc58bd63186bdeb9245.html'
	# PS学习：photoshop新手教程(珍藏版)  isPptlike = False, 与ppt不同，存在问题需解决
	pdf_url = 'https://wenku.baidu.com/view/1d496803bc64783e0912a21614791711cd797944.html?fr=search'
	# 人工智能的发展与应用。  图片并茂
	# doc_url = 'https://wenku.baidu.com/view/e9a90e449b6648d7c1c7469c.html?fr=search'
	# 淘宝宝贝模板制作教程
	# doc_url = 'https://wenku.baidu.com/view/9ff08107e87101f69e3195b0.html?fr=search'
	# 人工智能总结解读
	doc_url = 'https://wenku.baidu.com/view/fb92d7d3b8d528ea81c758f5f61fb7360a4c2b61.html?from=search'
	# python3基础教程
	# doc_url = 'https://wenku.baidu.com/view/b874f4db492fb4daa58da0116c175f0e7dd11977.html?fr=search'
	# 合同，82页
	# doc_url = 'https://wenku.baidu.com/view/61b1c212f18583d04964598c.html?fr=search'
	# 淘宝商品类目， 18页
	xls_url = 'https://wenku.baidu.com/view/a37d0f4fdc36a32d7375a417866fb84ae55cc360.html?fr=search'

	start = time.time()
	url = input('请输入百度文库链接：')
	doc = WenKuSpider(url, 'GetFiles')
	doc.run()
	end = time.time()
	print('耗时{}s'.format(str(end-start)))

	# 收费文档测试
	# ppt_url = 'https://wenku.baidu.com/view/e8229da7940590c69ec3d5bbfd0a79563c1ed47b'
	# ppt = WenKuSpider(ppt_url, 'GetFiles')
	# picurls = ppt.parse_ppt()

	# doc_url = 'https://wenku.baidu.com/view/359b8a37bed126fff705cc1755270722192e5978'
	# 图文
	# doc_url = 'https://wenku.baidu.com/view/588858feacaad1f34693daef5ef7ba0d4a736d3a'
	# doc = WenKuSpider(doc_url, 'GetFiles')
	# doc.run()


	# pdf_url = 'https://wenku.baidu.com/view/e0cb599fcec789eb172ded630b1c59eef9c79ad0'
	# pdf = WenKuSpider(pdf_url, 'GetFiles')
	# pdf.parse_ppt()



